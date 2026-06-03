#!/usr/bin/env python
"""Generate a presentation PPTX for the single-cell ANN search system.

Usage:
    python presentation/generate_ppt.py
    python presentation/generate_ppt.py --output presentation/单细胞ANN检索系统_中期展示.pptx

Requires: python-pptx (pip install python-pptx)
"""

from __future__ import annotations

import argparse
import os
import subprocess
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]

# ---------------------------------------------------------------------------
# Theme colours (matching front-end CSS)
# ---------------------------------------------------------------------------
GREEN = RGBColor(0x0F, 0x6F, 0x55)
WARM_WHITE = RGBColor(0xFF, 0xFD, 0xF7)
DARK_TEXT = RGBColor(0x11, 0x18, 0x14)
AMBER = RGBColor(0xB7, 0x79, 0x1F)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED = RGBColor(0x5F, 0x6D, 0x63)

CHARTS_DIR = PROJECT_ROOT / "benchmarks" / "charts"
OUTPUT_DIR = PROJECT_ROOT / "presentation"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def add_blank_slide(prs):
    """Add a blank slide with warm-white background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WARM_WHITE
    return slide


def add_title_bar(slide, title: str, subtitle: str = ""):
    """Top bar with green left accent and dark title."""
    # Green accent bar
    bar = slide.shapes.add_shape(
        1, Cm(0), Cm(2.5), Cm(25.4), Cm(2.2)  # rectangle, full width
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Cm(2), Cm(2.7), Cm(21), Cm(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WARM_WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xD0, 0xD8, 0xD2)
        p2.space_before = Pt(4)


def add_body_text(slide, left, top, width, height, lines: list[str], font_size=14):
    """Add a text box with bullet-point lines."""
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)
        if line.startswith("●") or line.startswith("•"):
            p.level = 1

    return txBox


def add_table(slide, left, top, headers: list[str], rows: list[list[str]], col_widths=None):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Cm(left), Cm(top), Cm(22), Cm(1 + n_rows * 0.8)
    )
    tbl = tbl_shape.table

    # Header
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WARM_WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN

    # Data
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WARM_WHITE if ri % 2 == 0 else RGBColor(0xF5, 0xF3, 0xEE)

    return tbl_shape


def add_image_safe(slide, path: Path, left, top, width=18, height=10):
    """Add an image if it exists, otherwise a placeholder message."""
    if path.exists():
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), Cm(width), Cm(height))
    else:
        txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"[图片缺失]\n{path.name}\n请先运行 benchmark_runner.py 和 benchmark_charts.py"
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.CENTER


def add_page_number(slide, num, total):
    """Add a small page number at the bottom-right."""
    txBox = slide.shapes.add_textbox(Cm(21), Cm(17.5), Cm(3), Cm(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_footer(slide):
    """Add a thin green line at the bottom."""
    line = slide.shapes.add_shape(1, Cm(0), Cm(18.5), Cm(25.4), Cm(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_cover(prs):
    """Slide 1: Title slide."""
    slide = add_blank_slide(prs)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = GREEN

    txBox = slide.shapes.add_textbox(Cm(3), Cm(5), Cm(20), Cm(8))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "单细胞 ANN 检索系统"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WARM_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "基于 FAISS 的高效单细胞 RNA 测序数据相似性检索"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xC8, 0xD5, 0xCC)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    p3 = tf.add_paragraph()
    p3.text = ""
    p3.space_before = Pt(36)

    p4 = tf.add_paragraph()
    p4.text = f"21st Group · 软件工程 · {datetime.now().strftime('%Y-%m-%d')}"
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(0xB0, 0xC0, 0xB4)
    p4.alignment = PP_ALIGN.CENTER


def slide_overview(prs, total):
    """Slide 2: Project overview."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "项目概述", "Project Overview")
    add_page_number(slide, 2, total)
    add_footer(slide)

    add_body_text(slide, 2, 5.5, 21, 12, [
        "问题背景",
        "● 单细胞 RNA 测序数据量大（数万细胞 × 数万基因），高维稀疏",
        "● 传统线性搜索在面对大规模数据时响应缓慢",
        "",
        "解决方案：ANN 近似最近邻检索",
        "● 使用 PCA 将基因表达降维到 30~50 维",
        "● 利用 FAISS 构建 ANN 索引加速 Top-K 相似细胞查询",
        "● 支持 Flat（精确）/ HNSW（图索引）/ IVF（倒排索引）三种索引",
        "● 提供 Web 界面，支持可视化（UMAP/PCA）和管理功能",
    ], font_size=14)


def slide_architecture(prs, total):
    """Slide 3: System architecture."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "系统架构", "System Architecture")
    add_page_number(slide, 3, total)
    add_footer(slide)

    add_body_text(slide, 2, 5.5, 21, 8, [
        "分层设计：",
        "",
        "┌─────────────────────────────────────────┐",
        "│  浏览器 (HTML + JS + Canvas)              │  前端",
        "└──────────────────┬──────────────────────┘",
        "                   ↓",
        "┌─────────────────────────────────────────┐",
        "│  Flask API (app.py)                       │  Web 层",
        "│  /api/search  /api/status  /api/rebuild  │",
        "└──────────────────┬──────────────────────┘",
        "                   ↓",
        "┌─────────────────────────────────────────┐",
        "│  SearchService (search_service.py)        │  业务层",
        "│  查询分发 · 结果格式化 · 索引切换          │",
        "└──────────────────┬──────────────────────┘",
        "                   ↓",
        "┌──────────────────────┬──────────────────┐",
        "│  ANNEngine            │  DataLoader       │  引擎层",
        "│  Flat / HNSW / IVF   │  PCA · 归一化     │",
        "│  (ann_engine.py)      │  (data_loader.py) │",
        "└──────────────────────┴──────────────────┘",
        "                   ↓                        ↓",
        "            FAISS                    .h5ad 数据",
    ], font_size=11)


def slide_pipeline(prs, total):
    """Slide 4: Data processing pipeline."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "数据处理流程", "Data Processing Pipeline")
    add_page_number(slide, 4, total)
    add_footer(slide)

    add_body_text(slide, 2, 5.5, 21, 11, [
        "原始 .h5ad 单细胞数据 → 预处理 → PCA 降维 → L2 归一化 → FAISS 索引",
        "",
        "1. 数据读取：load_h5ad() 读取 AnnData 对象",
        "   ● 检查 .h5ad 文件和 X_pca 是否已存在",
        "",
        "2. 预处理（如无 X_pca）：",
        "   ● normalize_total → log1p → highly_variable_genes → scale → PCA",
        "   ● 默认使用 Top-2000 高变基因，计算 50 个主成分",
        "",
        "3. L2 归一化：l2_normalize() 将所有向量归一化为单位向量",
        "   ● 内积 = 余弦相似度（FAISS 标准做法）",
        "",
        "4. 索引构建：ANNEngine.build_index()",
        "   ● Flat: 精确检索，O(n) 线性扫描",
        "   ● HNSW: 分层图，O(log n) 近似检索",
        "   ● IVF: 倒排索引，聚类中心 + 桶内搜索",
    ], font_size=13)


def slide_index_comparison(prs, total):
    """Slide 5: Three index types comparison table."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "三种索引对比", "Flat · HNSW · IVF")
    add_page_number(slide, 5, total)
    add_footer(slide)

    add_table(slide, 1.5, 5.5,
              ["特性", "Flat", "HNSW", "IVF"],
              [
                  ["算法类型", "精确暴力搜索", "分层可导航小世界图", "倒排文件索引"],
                  ["查询复杂度", "O(n)", "O(log n)", "O(√n)"],
                  ["是否需要训练", "否", "否", "是（需 k-means）"],
                  ["召回率", "1.0（精确）", "≈ 1.0", "≈ 1.0"],
                  ["适用场景", "小数据集 + 精度要求极高", "中等数据集 + 低延迟", "大数据集 + 内存敏感"],
                  ["我们的数据表现", "0.58 ms/查询", "0.09 ms/查询（6× 加速）", "0.14 ms/查询（4× 加速）"],
              ])


def slide_build_time(prs, total):
    """Slide 6: Build time performance."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "索引构建性能", "Build Time Comparison")
    add_page_number(slide, 6, total)
    add_footer(slide)

    img = CHARTS_DIR / "build_time_comparison.png"
    add_image_safe(slide, img, 2.5, 5.5, 21, 12)


def slide_query_latency(prs, total):
    """Slide 7: Query latency."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "查询延迟 (Top-K)", "Query Latency by Top-K")
    add_page_number(slide, 7, total)
    add_footer(slide)

    img = CHARTS_DIR / "query_latency_by_top_k.png"
    add_image_safe(slide, img, 2.5, 5.5, 21, 12)


def slide_recall(prs, total):
    """Slide 8: Recall vs latency trade-off."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "召回率与延迟权衡", "Recall vs Latency Trade-off")
    add_page_number(slide, 8, total)
    add_footer(slide)

    img = CHARTS_DIR / "latency_recall_tradeoff.png"
    add_image_safe(slide, img, 2.5, 5.5, 21, 12)


def slide_web_demo(prs, total):
    """Slide 9: Web interface features."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Web 界面功能", "Web Interface Features")
    add_page_number(slide, 9, total)
    add_footer(slide)

    add_body_text(slide, 2, 5.5, 10, 6, [
        "🔐 用户系统",
        "● 登录 / 注册",
        "● 管理员权限分级",
        "",
        "🔍 三种检索方式",
        "● 按细胞编号",
        "● 按真实 cell_id",
        "● 按自定义 PCA 向量",
    ], font_size=12)

    add_body_text(slide, 13, 5.5, 10, 6, [
        "📊 结果展示",
        "● Top-K 结果表格 + 耗时",
        "● metadata 筛选",
        "● UMAP/PCA canvas 可视化",
        "● 检索结果高亮标注",
        "",
        "⚙️ 管理功能",
        "● 数据集管理",
        "● 索引重建（Flat/HNSW/IVF）",
        "● 用户管理",
    ], font_size=12)

    # Screenshot placeholders
    add_body_text(slide, 2, 13, 21, 4, [
        "📸 截图清单见 presentation/screenshot_guide.md",
        "● 登录页 · 系统状态 · 检索结果 · metadata 筛选 · UMAP 可视化 · Top-K 高亮 · 管理后台 · 索引重建",
    ], font_size=11)


def slide_testing(prs, total):
    """Slide 10: Testing & quality."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "测试与质量保障", "Testing & Quality Assurance")
    add_page_number(slide, 10, total)
    add_footer(slide)

    add_body_text(slide, 2, 5.5, 10, 8, [
        "测试统计",
        "● 测试文件：8 个",
        "● 基础测试（不含真实数据）：90 个 ✅",
        "● 真实数据测试（liver.h5ad）：12 个 ✅",
        "● 总计：102 个测试，全部通过",
        "● 代码覆盖率：76%",
        "",
        "测试层次",
        "● 单元测试：核心模块独立验证",
        "● 集成测试：真实组件端到端流程",
        "● 真实数据测试：HNSW/IVF recall 验证",
        "● benchmark：性能基准评测",
    ], font_size=13)

    add_table(slide, 14, 5.5,
              ["模块", "覆盖率"],
              [
                  ["ann_engine.py", "90%"],
                  ["dataset_manager.py", "100%"],
                  ["db.py", "98%"],
                  ["search_service.py", "66%"],
                  ["user_service.py", "74%"],
                  ["data_loader.py", "40%"],
              ])


def slide_git(prs, total):
    """Slide 11: Project management."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "项目管理与协作", "Project Management")
    add_page_number(slide, 11, total)
    add_footer(slide)

    add_body_text(slide, 2, 5.5, 21, 10, [
        "Git 仓库结构",
        "● main 分支：主分支，合并所有功能",
        "● 各成员 feature 分支：yxy / gsd 等",
        "● 使用 Pull Request 合并",
        "",
        "模块分工（A / B / C / D）",
        "● A：数据加载与预处理（data_loader.py）",
        "● B：ANN 检索引擎（ann_engine.py）+ 索引管理",
        "● C：用户系统 + 数据集管理（SQLite）",
        "● D：Flask API + 前端 Web 界面 + 可视化",
        "",
        "协作工具",
        "● GitHub：代码托管、PR 审查",
        "● SQLite：轻量本地数据库，无需额外部署",
        "● pytest：自动化测试框架",
    ], font_size=13)


def slide_summary(prs, total):
    """Slide 12: Summary & next steps."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "总结与展望", "Summary & Future Work")
    add_page_number(slide, 12, total)
    add_footer(slide)

    add_body_text(slide, 2, 5.5, 21, 6, [
        "✅ 已完成",
        "● 完整的 ANN 检索系统：数据加载 → PCA → 索引构建 → 查询 → 可视化",
        "● Flat / HNSW / IVF 三种索引，全部集成并可运行时切换",
        "● 真实 liver.h5ad 数据集（69,032 细胞）验证通过",
        "● HNSW 查询延迟仅 0.09 ms，recall@10 = 1.0",
        "● 102 个自动化测试全部通过 + 性能 benchmark 体系",
        "",
        "📌 后续工作",
        "● 支持运行时动态切换数据集（/api/load-data）",
        "● 增加 FAISS GPU 支持，加速大规模数据",
        "● 增加批量查询接口",
        "● 增加 recall 调节参数（调大 efSearch/nprobe 换取更高召回）",
    ], font_size=13)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def get_git_stats() -> dict:
    """Collect git stats for the slide."""
    try:
        commits_raw = subprocess.check_output(
            ["git", "-C", str(PROJECT_ROOT), "log", "--oneline"],
            text=True, stderr=subprocess.DEVNULL,
        )
        commits = [l for l in commits_raw.strip().split("\n") if l]
        return {"n_commits": len(commits), "recent": commits[:5]}
    except Exception:
        return {"n_commits": "?", "recent": []}


def main():
    parser = argparse.ArgumentParser(description="Generate project PPTX")
    parser.add_argument(
        "--output",
        default=str(OUTPUT_DIR / "单细胞ANN检索系统_中期展示.pptx"),
        help="Output PPTX path",
    )
    args = parser.parse_args()

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    TOTAL_SLIDES = 12

    print("[ppt] Generating slides ...")
    slide_cover(prs)
    slide_overview(prs, TOTAL_SLIDES)
    slide_architecture(prs, TOTAL_SLIDES)
    slide_pipeline(prs, TOTAL_SLIDES)
    slide_index_comparison(prs, TOTAL_SLIDES)
    slide_build_time(prs, TOTAL_SLIDES)
    slide_query_latency(prs, TOTAL_SLIDES)
    slide_recall(prs, TOTAL_SLIDES)
    slide_web_demo(prs, TOTAL_SLIDES)
    slide_testing(prs, TOTAL_SLIDES)
    slide_git(prs, TOTAL_SLIDES)
    slide_summary(prs, TOTAL_SLIDES)

    prs.save(str(output_path))
    print(f"[ppt] Saved → {output_path}")
    print(f"[ppt] {TOTAL_SLIDES} slides generated.")


if __name__ == "__main__":
    main()
